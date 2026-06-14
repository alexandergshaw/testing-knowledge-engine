import io

import requests
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

body = {
    "title": "Foundations of Programming with Python",
    "objectives": (
        "Choose appropriate numeric data types. "
        "Organize computer programs using functions. "
        "Implement basic control structures."
    ),
}
r = requests.post("http://127.0.0.1:5050/api/v1/lecture", json=body, timeout=180)
r.raise_for_status()
prs = Presentation(io.BytesIO(r.content))

for i, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        paras = shape.text_frame.paragraphs
        if any(p.font.name == "Courier New" for p in paras):
            tf = shape.text_frame
            aligns = {str(p.alignment) for p in paras if p.text.strip()}
            print(f"Slide {i}: anchor={tf.vertical_anchor} aligns={aligns}")
            for p in paras:
                print("   |" + p.text)
