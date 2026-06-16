"""Deterministic text extraction from uploaded artifacts (decks, docs, sheets,
PDFs, plain text). Used to seed lecture generation from an existing file — the
extracted text is fed into the same objective-parsing + retrieval pipeline, so a
deck is just a richer way to supply `objectives`. No LLM, no OCR.

Office formats are ZIP+XML, so .pptx/.docx use the libraries already vendored
(python-pptx / python-docx) and .xlsx/.odf are read straight from the archive
with the standard library. PDF uses pypdf (lazy import; absent → no text). Every
extractor is best-effort: a corrupt or unreadable file yields "" rather than
raising, and the caller treats empty extraction as a 422.
"""

import io
import logging
import os
import re
import zipfile
from html import unescape

log = logging.getLogger(__name__)

# Plain-text / source files: decoded as UTF-8 directly.
_TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".rst", ".csv", ".tsv", ".json", ".yaml", ".yml",
    ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".h", ".cpp", ".cc",
    ".cs", ".go", ".rb", ".php", ".rs", ".swift", ".kt", ".sql", ".sh", ".html",
    ".css", ".xml",
}
_OFFICE_EXTS = {".pptx", ".docx", ".xlsx", ".odp", ".odt", ".ods"}
_OTHER_EXTS = {".pdf", ".rtf"}

SUPPORTED_EXTENSIONS = _TEXT_EXTS | _OFFICE_EXTS | _OTHER_EXTS


def _ext(filename):
    return os.path.splitext(filename or "")[1].lower()


def is_supported(filename):
    """True when the file's extension is one we can extract text from."""
    return _ext(filename) in SUPPORTED_EXTENSIONS


def _normalize(text):
    """Tidy extracted text: unify newlines, drop trailing spaces, collapse runs
    of blank lines, trim."""
    text = (text or "").replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _extract_pptx(data):
    from pptx import Presentation

    deck = Presentation(io.BytesIO(data))
    lines = []
    for slide in deck.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    lines.append(text)
    return "\n".join(lines)


def _extract_docx(data):
    from docx import Document

    document = Document(io.BytesIO(data))
    return "\n".join(p.text for p in document.paragraphs if p.text.strip())


def _extract_xlsx(data):
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        if "xl/sharedStrings.xml" not in archive.namelist():
            return ""
        xml = archive.read("xl/sharedStrings.xml").decode("utf-8", "replace")
    return "\n".join(unescape(t) for t in re.findall(r"<t[^>]*>(.*?)</t>", xml, re.S))


def _extract_odf(data):
    """OpenDocument (.odt/.odp/.ods): paragraphs live in content.xml."""
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        if "content.xml" not in archive.namelist():
            return ""
        xml = archive.read("content.xml").decode("utf-8", "replace")
    xml = re.sub(r"</text:(?:p|h|span)>", "\n", xml)
    xml = re.sub(r"<[^>]+>", "", xml)
    return unescape(xml)


def _extract_pdf(data):
    try:
        from pypdf import PdfReader
    except ImportError:
        log.warning("pypdf not installed; cannot extract PDF text")
        return ""
    reader = PdfReader(io.BytesIO(data))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _extract_rtf(data):
    """Best-effort RTF: drop control words and groups, keep the prose."""
    text = data.decode("latin-1", "replace")
    text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)   # hex-escaped chars
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", text)  # control words
    text = re.sub(r"[{}]", "", text)
    return text


def _extract_text(data):
    return data.decode("utf-8", "replace")


_EXTRACTORS = {
    ".pptx": _extract_pptx,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".odp": _extract_odf,
    ".odt": _extract_odf,
    ".ods": _extract_odf,
    ".pdf": _extract_pdf,
    ".rtf": _extract_rtf,
}


def extract_text(filename, data):
    """Plain text from an uploaded file's bytes, or "" when nothing is
    extractable (empty, scanned, corrupt, or an unsupported type). Never raises."""
    ext = _ext(filename)
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        extractor = _extract_text if ext in _TEXT_EXTS else None
    if extractor is None:
        return ""
    try:
        return _normalize(extractor(data))
    except Exception:
        log.warning("text extraction failed for %s", filename, exc_info=True)
        return ""


# --- outline (headings/titles only) -----------------------------------------
# For seeding a lecture from a deck/doc, the slide titles and headings are the
# topic signal; body prose and code are noise. extract_outline returns just that.


def _looks_like_heading(line):
    line = line.strip()
    if not 2 <= len(line) <= 90:
        return False
    if line.endswith((".", "!", "?", ":", ";", ",")):  # a sentence/clause, not a heading
        return False
    if re.search(r"==|!=|->|::|[(){}\[\]<>]|[a-zA-Z_]\w*\s*=\s*\S", line):  # code-ish
        return False
    if sum(ch.isdigit() for ch in line) > len(line) * 0.4:  # mostly numbers
        return False
    return True


def _outline_pptx(data):
    from pptx import Presentation

    deck = Presentation(io.BytesIO(data))
    titles = []
    for slide in deck.slides:
        title = ""
        if slide.shapes.title is not None and slide.shapes.title.text.strip():
            title = slide.shapes.title.text.strip()
        else:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    title = shape.text_frame.paragraphs[0].text.strip()
                    break
        if title:
            titles.append(title.splitlines()[0].strip())
    return titles


def _outline_docx(data):
    from docx import Document

    document = Document(io.BytesIO(data))
    out = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = (paragraph.style.name or "").lower() if paragraph.style else ""
        if style.startswith("heading") or style == "title" or _looks_like_heading(text):
            out.append(text)
    return out


def extract_outline(filename, data):
    """Headings/slide-titles from a file, one per line, deduped — the clean topic
    signal. Falls back to heading-shaped lines from the full text for formats
    without explicit headings. Returns "" on failure."""
    ext = _ext(filename)
    try:
        if ext == ".pptx":
            lines = _outline_pptx(data)
        elif ext == ".docx":
            lines = _outline_docx(data)
        else:
            lines = [ln for ln in extract_text(filename, data).split("\n") if _looks_like_heading(ln)]
    except Exception:
        log.warning("outline extraction failed for %s", filename, exc_info=True)
        lines = []
    seen, out = set(), []
    for line in lines:
        line = line.strip()
        if line and line not in seen:
            seen.add(line)
            out.append(line)
    return "\n".join(out)
