"""Shared python-pptx slide builders + the lecture theme, mimicking the
reference "Gemini" deck: a navy header band with a white title, a light
background, bright-blue accents, and dark code blocks with a language label.

House rules: 16:9, Title-Cased self-contained titles, up to six short
self-contained bullets per content slide (overflow → speaker notes), compact
non-bullet list slides for agendas/references."""

import io
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

MAX_BULLETS = 6
SLIDE_W = 13.333
SLIDE_H = 7.5
SLIDE_WIDTH = Inches(SLIDE_W)
SLIDE_HEIGHT = Inches(SLIDE_H)
HEADER_H = 1.4
BLANK_LAYOUT = 6


class Theme:
    title_font = "Calibri"
    body_font = "Calibri"
    mono_font = "Courier New"
    header = RGBColor(0x1A, 0x27, 0x44)        # navy header band / title slide bg
    background = RGBColor(0xF4, 0xF6, 0xFB)     # light content background
    accent = RGBColor(0x25, 0x63, 0xEB)        # bright blue
    body = RGBColor(0x1E, 0x29, 0x3B)          # dark slate body text
    on_header = RGBColor(0xFF, 0xFF, 0xFF)      # white title text
    muted = RGBColor(0x64, 0x74, 0x8B)
    code_bg = RGBColor(0x0F, 0x17, 0x2A)        # dark code block
    code_text = RGBColor(0xE2, 0xE8, 0xF0)      # light code text


THEME = Theme()

_SMALL_WORDS = {
    "a", "an", "the", "and", "or", "but", "nor", "of", "to", "in", "on", "at",
    "for", "with", "by", "as", "is", "are", "vs", "via", "per",
}


# --- text normalization ------------------------------------------------------


def clean_title(text):
    """Trim, drop trailing punctuation, Title-Case while keeping acronyms /
    mixed-case tokens (Python, I/O, HTML) intact."""
    text = re.sub(r"\s+", " ", str(text)).strip().rstrip(".,;:!")
    words = text.split(" ")
    out = []
    for index, word in enumerate(words):
        edge = index == 0 or index == len(words) - 1
        if not word:
            continue
        if (len(word) > 1 and any(c.isupper() for c in word[1:])) or (word.isupper() and len(word) > 1):
            out.append(word)
        elif not edge and word.lower() in _SMALL_WORDS:
            out.append(word.lower())
        else:
            # Capitalize across hyphens: "problem-solving" -> "Problem-Solving".
            out.append("-".join(part[:1].upper() + part[1:].lower() for part in word.split("-")))
    return " ".join(out)[:90]


def clean_bullet(text):
    """Make a bullet stand alone: strip leading list/citation markers,
    capitalize, ensure terminal punctuation."""
    text = re.sub(r"\s+", " ", str(text)).strip()
    text = re.sub(r"^(?:[-*•–]|\d+[.)])\s*", "", text)
    text = re.sub(r"\s*\[\d+\]", "", text)
    text = text.strip()
    if not text:
        return ""
    text = text[:1].upper() + text[1:]
    if text[-1] not in ".!?:":
        text += "."
    return text


# --- low-level shape helpers -------------------------------------------------


def _blank(deck):
    return deck.slides.add_slide(deck.slide_layouts[BLANK_LAYOUT])


def _fill_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _text(slide, text, x, y, w, h, size, color, bold=False, font=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0)
    frame.margin_right = Inches(0)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = font or THEME.body_font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def _content_base(deck, title):
    """A content slide: light background, navy header band with white title, a
    thin blue accent line under it and a short blue left rail (Gemini style)."""
    slide = _blank(deck)
    _fill_background(slide, THEME.background)
    _rect(slide, 0, 0, SLIDE_W, HEADER_H, THEME.header)
    _rect(slide, 0, HEADER_H, SLIDE_W, 0.06, THEME.accent)
    _rect(slide, 0, HEADER_H, 0.08, SLIDE_H - HEADER_H, THEME.accent)
    _text(slide, clean_title(title), 0.45, 0.28, 12.4, 0.95, 26, THEME.on_header, bold=True,
          font=THEME.title_font)
    return slide


# --- public builders ---------------------------------------------------------


def new_deck():
    deck = Presentation()
    deck.slide_width = SLIDE_WIDTH
    deck.slide_height = SLIDE_HEIGHT
    return deck


def add_title_slide(deck, title, subtitle=""):
    slide = _blank(deck)
    _fill_background(slide, THEME.header)
    _rect(slide, 0, 4.6, SLIDE_W, 0.10, THEME.accent)   # horizontal accent
    _rect(slide, 0, 0, 0.20, SLIDE_H, THEME.accent)     # left accent rail
    _text(slide, clean_title(title), 0.7, 2.2, 12.0, 2.0, 42, THEME.on_header, bold=True,
          font=THEME.title_font)
    if subtitle:
        _text(slide, subtitle, 0.7, 4.9, 12.0, 0.8, 20, RGBColor(0xCA, 0xDC, 0xFC),
              font=THEME.body_font)
    return slide


def _body_paragraphs(slide, lines, top=1.75, size=18, bullet=True):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(12.3), Inches(SLIDE_H - top - 0.4))
    frame = box.text_frame
    frame.word_wrap = True
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = ("•  " + line) if bullet else line
        paragraph.font.name = THEME.body_font
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = THEME.body
        paragraph.space_after = Pt(10)
    return box


def add_bullet_slide(deck, heading, bullets, notes="", footer=None):
    """A content slide with up to MAX_BULLETS short, self-contained bullets;
    any overflow (plus extra notes) goes to the slide's speaker notes."""
    slide = _content_base(deck, heading)
    cleaned = [b for b in (clean_bullet(item) for item in bullets) if b]
    _body_paragraphs(slide, cleaned[:MAX_BULLETS])

    note_parts = []
    overflow = cleaned[MAX_BULLETS:]
    if overflow:
        note_parts.append("Also covers:\n" + "\n".join(f"- {b}" for b in overflow))
    if notes:
        note_parts.append(notes)
    if note_parts:
        set_notes(slide, "\n\n".join(note_parts))
    return slide


def add_list_slide(deck, heading, items, footer=None):
    """Compact non-bullet list (small font) for agendas and references."""
    slide = _content_base(deck, heading)
    box = slide.shapes.add_textbox(Inches(0.55), Inches(1.75), Inches(12.3), Inches(5.1))
    frame = box.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = str(item)
        paragraph.font.name = THEME.body_font
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = THEME.body
        paragraph.space_after = Pt(6)
    return slide


def add_content_slide(deck, title, footer=None):
    """A themed content slide we populate with custom text/code boxes."""
    return _content_base(deck, title)


def add_text_box(slide, text, top=1.75, height=1.4, size=18):
    return _text(slide, text, 0.55, top, 12.3, height, size, THEME.body, font=THEME.body_font)


def add_code_box(slide, lines, top=3.5, height=3.3):
    """A dark code block with light monospace text (Gemini style). Auto-shapes
    default to centered, vertically-centered text — fatal for code — so text is
    forced left-aligned and top-anchored, with indentation preserved."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(top), Inches(12.3), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = THEME.code_bg
    box.line.fill.background()
    box.shadow.inherit = False
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.margin_left = Inches(0.25)
    frame.margin_top = Inches(0.18)
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line.replace("\t", "    ")
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.name = THEME.mono_font
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = THEME.code_text
    return box


def add_code_example_slide(deck, title, caption, language, lines):
    """A code example slide: caption, a blue language label, then a dark code
    block — the reference deck's 'Example: …' layout."""
    slide = _content_base(deck, title)
    if caption:
        add_text_box(slide, caption, top=1.7, height=1.2, size=18)
    if language:
        _text(slide, language.upper(), 0.55, 3.15, 12.3, 0.3, 11, THEME.accent, bold=True,
              font=THEME.title_font)
    add_code_box(slide, lines, top=3.5, height=3.3)
    return slide


def slide_title(slide):
    """Title text of a slide built by these helpers (the topmost text box)."""
    titled = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
    if not titled:
        return ""
    return min(titled, key=lambda s: s.top or 0).text_frame.text.strip()


def slide_bullets(slide):
    """The bullet lines on a content slide (the box of '•  '-prefixed lines)."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        paragraphs = [p.text for p in shape.text_frame.paragraphs if p.text.strip()]
        if paragraphs and all(p.lstrip().startswith("•") for p in paragraphs):
            return paragraphs
    return []


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def deck_bytes(deck):
    output = io.BytesIO()
    deck.save(output)
    return output.getvalue()
