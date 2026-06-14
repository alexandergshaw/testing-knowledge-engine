import io

from pptx import Presentation

from knowledge.lecture import ObjectiveResult, build_module_deck
from knowledge.slides import (
    MAX_BULLETS,
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
    add_bullet_slide,
    clean_bullet,
    clean_title,
    new_deck,
    slide_bullets,
    slide_title,
)


# --- text normalization ------------------------------------------------------


def test_clean_title_title_cases_and_preserves_acronyms():
    assert (
        clean_title("objective 1: explain what a python variable is")
        == "Objective 1: Explain What a Python Variable Is"
    )
    assert clean_title("intro to HTML and CSS") == "Intro to HTML and CSS"
    assert clean_title("Variables, I/O, Branching.") == "Variables, I/O, Branching"
    assert clean_title("problem-solving strategies") == "Problem-Solving Strategies"


def test_clean_bullet_makes_self_contained():
    assert clean_bullet("[1] the mental discomfort felt") == "The mental discomfort felt."
    assert clean_bullet("- run the tests!") == "Run the tests!"
    assert clean_bullet("   ") == ""


# --- slide invariants --------------------------------------------------------


def test_code_box_left_aligned_and_top_anchored():
    # Auto-shapes default to centered text — code must be forced left + top.
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    from knowledge.slides import add_code_box, add_content_slide

    deck = new_deck()
    slide = add_content_slide(deck, "Example")
    box = add_code_box(slide, ["def f():", "    return 1"])
    assert box.text_frame.vertical_anchor == MSO_ANCHOR.TOP
    for paragraph in box.text_frame.paragraphs:
        if paragraph.text.strip():
            assert paragraph.alignment == PP_ALIGN.LEFT


def test_new_deck_is_widescreen():
    deck = new_deck()
    assert deck.slide_width == SLIDE_WIDTH
    assert deck.slide_height == SLIDE_HEIGHT


def test_bullet_slide_caps_and_overflows_to_notes():
    deck = new_deck()
    add_bullet_slide(deck, "Heading", [f"Point {i} stands on its own." for i in range(8)])
    slide = deck.slides[0]
    shown = slide_bullets(slide)
    assert len(shown) == MAX_BULLETS == 6

    notes = slide.notes_slide.notes_text_frame.text
    assert "Point 6 stands on its own." in notes  # 7th point overflowed


def _sample_results():
    return [
        ObjectiveResult(
            objective="define variables",
            points=[f"Point {i} is a complete, self-contained sentence." for i in range(8)],
            examples=[
                {"kind": "prose", "text": "For example, x = 5.", "title": "V", "url": "u", "source": "Wikipedia"}
            ],
            citations=[{"title": "V", "url": "u", "source": "Wikipedia"}],
            confidence="high",
        ),
        ObjectiveResult(
            objective="organize programs",  # names no concept -> explanation only
            points=["A program is an organized set of steps."],
            examples=[],
            citations=[],
            confidence="medium",
        ),
    ]


def test_module_deck_content_slides_obey_bullet_cap():
    deck = Presentation(io.BytesIO(build_module_deck("Intro to Python", _sample_results())))
    assert deck.slide_width == SLIDE_WIDTH  # widescreen + themed
    counts = [len(slide_bullets(s)) for s in deck.slides if slide_bullets(s)]
    assert counts  # at least one content slide has bullets
    assert all(count <= MAX_BULLETS for count in counts)


def test_module_deck_overflow_points_land_in_notes():
    deck = Presentation(io.BytesIO(build_module_deck("Intro to Python", _sample_results())))
    objective_slide = next(
        s for s in deck.slides if slide_title(s).startswith("Variables")
    )
    notes = objective_slide.notes_slide.notes_text_frame.text
    # 8 points, only 6 shown -> the rest are preserved in the notes.
    assert "Point 7 is a complete" in notes
