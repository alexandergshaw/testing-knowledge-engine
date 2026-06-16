import io
import zipfile

from docx import Document
from pptx import Presentation
from pptx.util import Inches

from knowledge import extract


def make_pptx(texts):
    deck = Presentation()
    blank = deck.slide_layouts[6]
    for text in texts:
        slide = deck.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        box.text_frame.text = text
    buf = io.BytesIO()
    deck.save(buf)
    return buf.getvalue()


def make_docx(paragraphs):
    document = Document()
    for para in paragraphs:
        document.add_paragraph(para)
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def make_xlsx(strings):
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as archive:
        items = "".join(f"<si><t>{s}</t></si>" for s in strings)
        archive.writestr("xl/sharedStrings.xml", f'<?xml version="1.0"?><sst>{items}</sst>')
    return buf.getvalue()


def test_is_supported():
    assert extract.is_supported("deck.pptx")
    assert extract.is_supported("notes.MD")           # case-insensitive
    assert extract.is_supported("report.pdf")
    assert not extract.is_supported("malware.exe")
    assert not extract.is_supported("legacy.ppt")     # binary legacy formats unsupported


def test_extract_plain_text():
    assert extract.extract_text("a.txt", b"Define variables\nExplain loops") == "Define variables\nExplain loops"


def test_extract_pptx_concatenates_slide_text():
    data = make_pptx(["Introduction to Photosynthesis", "Light reactions", "The Calvin cycle"])
    text = extract.extract_text("deck.pptx", data)
    assert "Introduction to Photosynthesis" in text
    assert "Calvin cycle" in text


def test_extract_docx_paragraphs():
    data = make_docx(["Newton's first law", "Newton's second law", "Newton's third law"])
    text = extract.extract_text("doc.docx", data)
    assert "Newton's first law" in text and "third law" in text


def test_extract_xlsx_shared_strings():
    data = make_xlsx(["Supply", "Demand", "Equilibrium"])
    text = extract.extract_text("sheet.xlsx", data)
    assert "Supply" in text and "Equilibrium" in text


def test_corrupt_file_returns_empty_not_raises():
    assert extract.extract_text("deck.pptx", b"not a real zip") == ""


def test_unsupported_extension_returns_empty():
    assert extract.extract_text("x.exe", b"\x00\x01") == ""


def test_normalize_collapses_blank_lines():
    assert extract.extract_text("a.txt", b"A\n\n\n\nB\n  \n  C  ") == "A\n\nB\n\n  C"


def test_extract_outline_pptx_returns_titles():
    data = make_pptx(["Introduction", "Boolean Expressions", "While Loop"])
    assert extract.extract_outline("ch.pptx", data).split("\n") == [
        "Introduction", "Boolean Expressions", "While Loop"
    ]


def test_extract_outline_generic_drops_prose_and_code():
    text = b"Boolean Expressions\nThis sentence explains things in great detail.\na == b and c\nWhile Loop"
    outline = extract.extract_outline("ch.txt", text).split("\n")
    assert "Boolean Expressions" in outline and "While Loop" in outline
    assert not any("==" in line for line in outline)          # code dropped
    assert not any(line.endswith(".") for line in outline)    # prose sentence dropped
