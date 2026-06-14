import io

from pptx import Presentation

from knowledge import case_study
from knowledge.lecture import ObjectiveResult, build_module_deck
from knowledge.slides import slide_title


def test_every_case_study_is_well_formed():
    entries = list(case_study.CASE_STUDIES.values()) + [case_study.DEFAULT]
    for entry in entries:
        assert entry["title"].startswith("Case Study:")
        assert 1 <= len(entry["bullets"]) <= 4          # lecture-deck bullet budget
        assert entry["source"]["url"].startswith("https://")
        assert "code" not in entry and "lines" not in entry  # never carries code


def test_domain_matching_picks_the_specific_study():
    assert case_study.case_study_for("Introduction to Cybersecurity") is case_study.CASE_STUDIES["security"]
    assert case_study.case_study_for("Machine Learning 101") is case_study.CASE_STUDIES["ml_ai"]
    assert case_study.case_study_for("Relational Databases", "write SQL queries") is case_study.CASE_STUDIES["databases"]
    assert case_study.case_study_for("Operating Systems", "threads and deadlock") is case_study.CASE_STUDIES["os_concurrency"]
    assert case_study.case_study_for("Web Development", "build a website with HTML") is case_study.CASE_STUDIES["web"]
    assert case_study.case_study_for("Algorithms", "sorting and complexity") is case_study.CASE_STUDIES["algorithms"]


def test_unmatched_subject_falls_back_to_default():
    # A generic intro module names no domain — still gets a real case study.
    assert case_study.case_study_for("Introduction to Python") is case_study.DEFAULT
    assert case_study.case_study_for("Foundations of Computer Science") is case_study.DEFAULT


def test_case_study_slide_is_slide_three_with_no_code():
    results = [ObjectiveResult(objective="Define variables", points=["A variable stores data."],
                               confidence="high")]
    deck = Presentation(io.BytesIO(build_module_deck("Introduction to Cybersecurity", results)))
    titles = [slide_title(s) for s in deck.slides]
    assert titles[2].startswith("Case Study:")          # after title + overview
    assert sum(t.startswith("Case Study:") for t in titles) == 1

    slide = deck.slides[2]
    # An ordinary bullets slide: no monospace code box anywhere on it.
    assert not any(
        shape.has_text_frame and any(p.font.name == "Courier New" for p in shape.text_frame.paragraphs)
        for shape in slide.shapes
    )
    # The real source is cited in the speaker notes (deterministic, checkable).
    assert "wikipedia.org" in slide.notes_slide.notes_text_frame.text
