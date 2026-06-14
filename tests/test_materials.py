import csv
import io
import json
import zipfile

import pytest
from docx import Document
from pptx import Presentation

from knowledge.materials import (
    CSV_COLUMNS,
    MaterialsError,
    build_materials,
    build_rubric,
    parse_project,
    rubric_to_csv,
)

INSTRUCTIONS = """# assignment1: Variables, I/O, Branching

## Starter workflow (GUI-only)
- Edit only `student_code.py`.
- Run `test_assignment.py` from the Testing panel.

## Learning target
Build your daily summary logic so the dashboard widget unlocks.

## Worked example (different data)
Worked example data: temperatures = [68, 71, 75].
"""

STUDENT_CODE = '''\
# =============================================================================
#  ASSIGNMENT 1 — Variables, I/O, and Branching
# =============================================================================
#
#  WHAT YOU'RE LEARNING THIS WEEK
#
#  1. VARIABLES — A variable is a named storage container for a piece of data.
#     You create one by writing:  name = value
#
#       temperature = 72          # an integer
#       city = "Denver"           # a string
#
#  2. BRANCHING (if / elif / else) — Branching lets your program make decisions.
#
#       if score >= 90:
#           grade = "A"
# =============================================================================

"""Starter code for assignment1."""

student_name = "Your Name"


def get_dashboard_payload():
    """Return dashboard data.

    Your job:
      1. Create some variables to represent your daily data.
      2. Build labels and values lists from your data.

    Remember the rules:
      - At least 3 items in each list.
      - Same number of labels and values.
    """
    # ── Step 1: Define your data as variables ──────────
    # ── Step 2: Assemble your lists ────────────────────
    my_labels = []
    my_values = []
    return {"title": "x", "values": my_values, "labels": my_labels}
'''

TEST_CODE = """\
import unittest


class StudentCodeContractTests(unittest.TestCase):
    def test_get_dashboard_payload_exists(self):
        pass

    def test_get_dashboard_payload_shape(self):
        pass
"""

FINAL_INSTRUCTIONS = """# final: Final Project Integration

## Learning target
Integrate every prior week's work into the capstone.
"""


def make_project_zip():
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr("proj-main/README.md", "# proj")
        archive.writestr("proj-main/assignments/assignment1/INSTRUCTIONS.md", INSTRUCTIONS)
        archive.writestr("proj-main/assignments/assignment1/student_code.py", STUDENT_CODE)
        archive.writestr("proj-main/assignments/assignment1/test_assignment.py", TEST_CODE)
        archive.writestr("proj-main/assignments/final/INSTRUCTIONS.md", FINAL_INSTRUCTIONS)
        archive.writestr("proj-main/assignments/final/student_code.py", "x = 1\n")
        archive.writestr(
            "proj-main/assignments/assignment1/__pycache__/junk.pyc", "ignore"
        )
    return buffer.getvalue()


def test_parse_project_extracts_units():
    units = parse_project(make_project_zip())
    assert [u.slug for u in units] == ["assignment1", "final"]
    unit = units[0]
    assert unit.week == 1
    assert unit.topic == "Variables, I/O, Branching"
    assert [c["name"] for c in unit.concepts] == [
        "Variables",
        "Branching (if / elif / else)",
    ]
    assert 'temperature = 72          # an integer' in unit.concepts[0]["example"]
    assert unit.tasks[0].startswith("Create some variables")
    assert unit.rules[0].startswith("At least 3 items")
    assert unit.steps == ["Define your data as variables", "Assemble your lists"]
    assert unit.tests == [
        "test_get_dashboard_payload_exists",
        "test_get_dashboard_payload_shape",
    ]
    assert 'student_name = "Your Name"' in unit.placeholder_lines
    assert "my_labels = []" in unit.placeholder_lines


def test_rubric_is_deterministic_and_weights_sum_to_100():
    units = parse_project(make_project_zip())
    rubric_a = build_rubric(units)
    rubric_b = build_rubric(parse_project(make_project_zip()))
    assert json.dumps(rubric_a) == json.dumps(rubric_b)  # deterministic
    assert rubric_a["version"] == 1
    assert set(rubric_a["grader_contract"]) >= {
        "pytest",
        "forbidden_lines_absent",
        "python_compiles",
        "files_present",
    }
    for unit in rubric_a["units"]:
        assert sum(c["weight"] for c in unit["criteria"]) == 100
        assert unit["max_points"] == 100
    first = rubric_a["units"][0]["criteria"][0]
    assert first["type"] == "pytest"
    assert first["tests"] == [
        "test_get_dashboard_payload_exists",
        "test_get_dashboard_payload_shape",
    ]


def _docx_text(blob):
    document = Document(io.BytesIO(blob))
    return [(p.style.name, p.text) for p in document.paragraphs if p.text]


def test_build_materials_bundle_contents():
    payload, summary = build_materials(make_project_zip())
    assert summary["units"] == 2
    with zipfile.ZipFile(io.BytesIO(payload)) as bundle:
        names = bundle.namelist()
        assert "rubric.csv" in names
        assert "MANIFEST.docx" in names
        assert not any(n.endswith((".md", ".json")) for n in names)
        lectures = [n for n in names if n.startswith("lectures/") and n.endswith(".pptx")]
        lms = [n for n in names if n.startswith("lms/") and n.endswith(".docx")]
        assignments = [
            n for n in names if n.startswith("assignments/") and n.endswith(".docx")
        ]
        assert len(lectures) == len(lms) == len(assignments) == 2

        deck = Presentation(io.BytesIO(bundle.read(lectures[0])))
        assert len(deck.slides) >= 5  # title, goals, 2 concepts, task/grading
        # Every content slide carries at most MAX_BULLETS bullets.
        from knowledge.slides import MAX_BULLETS, slide_bullets

        for slide in deck.slides:
            assert len(slide_bullets(slide)) <= MAX_BULLETS

        intro = _docx_text(bundle.read([n for n in lms if "week-01" in n][0]))
        assert ("Heading 1", "Week 1: Variables, I/O, Branching") in intro
        assert ("Heading 2", "How you'll be graded") in intro
        assert any(style == "List Bullet" for style, _ in intro)

        doc = _docx_text(bundle.read([n for n in assignments if "week-01" in n][0]))
        assert ("Heading 1", "Week 1 Assignment — Variables, I/O, Branching") in doc
        assert ("Heading 2", "Your job") in doc
        assert ("Heading 2", "Rules your solution must follow") in doc
        assert any(style == "List Number" for style, _ in doc)

        manifest = Document(io.BytesIO(bundle.read("MANIFEST.docx")))
        assert len(manifest.tables) == 1
        assert len(manifest.tables[0].rows) == 3  # header + 2 units


def test_rubric_csv_shape():
    units = parse_project(make_project_zip())
    text = rubric_to_csv(build_rubric(units))
    assert text == rubric_to_csv(build_rubric(units))  # deterministic
    rows = list(csv.DictReader(io.StringIO(text)))
    assert list(rows[0]) == CSV_COLUMNS

    contract_rows = [r for r in rows if r["unit_id"] == "GRADER_CONTRACT"]
    assert {r["criterion_type"] for r in contract_rows} >= {
        "pytest",
        "forbidden_lines_absent",
        "python_compiles",
        "files_present",
    }

    unit_rows = [r for r in rows if r["unit_id"] == "assignment1"]
    assert sum(int(r["weight"]) for r in unit_rows) == 100
    pytest_row = next(r for r in unit_rows if r["criterion_type"] == "pytest")
    assert "test_get_dashboard_payload_exists" in pytest_row["details"]
    placeholder_row = next(
        r for r in unit_rows if r["criterion_type"] == "forbidden_lines_absent"
    )
    assert 'student_name = "Your Name"' in placeholder_row["details"]


def test_invalid_zip_rejected():
    with pytest.raises(MaterialsError):
        parse_project(b"this is not a zip")


def test_zip_without_units_rejected():
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr("readme.txt", "nothing here")
    with pytest.raises(MaterialsError):
        parse_project(buffer.getvalue())
