import io

import pytest
from pptx import Presentation

import app as app_module
import service as service_module
from knowledge.lecture import (
    ObjectiveResult,
    build_module_deck,
    extract_examples,
    parse_objectives,
)
from knowledge.query import analyze
from knowledge.ranking import ScoredSentence
from knowledge.slides import slide_title
from knowledge.sources.base import Passage
from service import PPTX_MIMETYPE


# --- parse_objectives: format-agnostic --------------------------------------


def test_parse_lead_in_prose():
    objectives = parse_objectives(
        "By the end of this module, students will be able to define variables, "
        "explain control flow, and write functions."
    )
    assert objectives == ["define variables", "explain control flow", "write functions"]


def test_parse_inline_numbered_list():
    objectives = parse_objectives("1. Define variables 2. Explain loops 3. Write functions")
    assert objectives == ["Define variables", "Explain loops", "Write functions"]


def test_parse_inline_bullets():
    objectives = parse_objectives("- Define variables - Explain loops - Write functions")
    assert objectives == ["Define variables", "Explain loops", "Write functions"]


def test_parse_run_on_action_verbs():
    objectives = parse_objectives(
        "Define variables. Explain control flow. Apply reinforcement schedules."
    )
    assert objectives == [
        "Define variables",
        "Explain control flow",
        "Apply reinforcement schedules",
    ]


def test_parse_single_objective():
    text = "Understand the fundamentals of cognitive dissonance theory"
    assert parse_objectives(text) == [text]


def test_parse_does_not_truncate_outcomes_objective():
    # "outcomes" is a lead-in word only as a colon header — not mid-objective.
    objectives = parse_objectives(
        "Evaluate learning outcomes; design assessment rubrics"
    )
    assert objectives == ["Evaluate learning outcomes", "design assessment rubrics"]


def test_parse_caps_and_dedupes():
    objectives = parse_objectives("define x; define x; explain y")
    assert objectives == ["define x", "explain y"]


# --- extract_examples --------------------------------------------------------


def make_passage(text="An explanation.", title="T", source="Wikipedia", code=None):
    return Passage(text=text, title=title, url="http://x", source=source, trust=1.0, code=code or [])


def test_extract_code_example_for_programming():
    query = analyze("How do Python loops work?")
    assert query.is_programming
    passage = make_passage(source="Stack Overflow", code=["for i in range(3):\n    print(i)"])
    examples = extract_examples(query, [passage], [])
    assert examples[0]["kind"] == "code"
    assert examples[0]["lines"][0] == "for i in range(3):"
    assert examples[0]["lines"][1] == "    print(i)"


def test_programming_lecture_defers_to_concept_examples():
    # In a programming lecture, per-objective examples are empty — concepts drive
    # the example slides instead (see _attach_concept_examples).
    query = analyze("How do Python loops work?")
    passage = make_passage(source="Stack Overflow", code=["for i in range(3):\n    print(i)"])
    assert extract_examples(query, [passage], [], programming_lecture=True) == []


# --- concept extraction & per-concept code ----------------------------------


def test_extract_concepts_ordered_and_deduped():
    from knowledge.lecture import extract_concepts

    concepts = extract_concepts(analyze("Explain conditionals and loops, and write functions"))
    assert concepts == ["Conditionals", "Loops", "Functions"]
    # synonyms collapse to the canonical name
    assert extract_concepts(analyze("Use lists and arrays")) == ["Lists & Arrays"]
    # a conceptual / language-only objective names no concept
    assert extract_concepts(analyze("Explain the history of Python")) == []


def test_extract_concepts_multiword_phrases():
    from knowledge.lecture import extract_concepts

    assert extract_concepts(analyze("Implement basic control structures")) == ["Control Structures"]
    assert extract_concepts(analyze("Choose appropriate numeric data types")) == ["Data Types"]
    assert extract_concepts(analyze("Explain common data structures")) == ["Data Structures"]
    # genuinely conceptual objectives still name no concept (no code slide)
    assert extract_concepts(analyze("Describe problem-solving strategies")) == []
    assert extract_concepts(analyze("Provide examples of computer science")) == []


def test_parse_recognizes_choose_and_provide_verbs():
    # "Choose"/"Provide" must start their own objectives, not merge into the prior.
    objs = parse_objectives(
        "Provide examples of CS. Choose numeric data types. Implement control structures."
    )
    assert objs == [
        "Provide examples of CS",
        "Choose numeric data types",
        "Implement control structures",
    ]


def test_concept_code_forces_stackoverflow_retrieval(monkeypatch):
    # A concept search string ("Data Types example") doesn't trip is_programming,
    # so _concept_code must force it so Stack Overflow (the code source) is queried.
    import knowledge.lecture as lecture

    captured = {}

    def fake_select(query):
        captured["is_programming"] = query.is_programming
        return []

    monkeypatch.setattr(lecture, "select_sources", fake_select)
    monkeypatch.setattr(lecture, "fetch", lambda query, sources: [])
    lecture._concept_code("Data Types", "uniquelang-xyz")  # fresh cache key
    assert captured["is_programming"] is True


def test_infer_language_from_title_and_objectives():
    from knowledge.lecture import _infer_language

    assert _infer_language(["write a for loop"], "Introduction to Python") == "Python"
    assert _infer_language(["write JavaScript functions"], "Web Dev") == "JavaScript"
    assert _infer_language(["explain cognitive dissonance"], "Psychology") == ""


def test_concept_code_caches_and_captions(monkeypatch):
    import knowledge.lecture as lecture

    calls = {"n": 0}

    def fake_fetch(query, sources):
        calls["n"] += 1
        return [make_passage(source="Stack Overflow", code=["for i in range(3):\n    print(i)"])]

    monkeypatch.setattr(lecture, "fetch", fake_fetch)
    monkeypatch.setattr(lecture, "select_sources", lambda q: [])
    monkeypatch.setattr(
        lecture, "rank",
        lambda q, passages: [
            ScoredSentence(text="For example, a loop repeats a block.", tokens=[], passage=passages[0], score=5.0)
        ],
    )

    example = lecture._concept_code("Loops", "Python")
    assert example["kind"] == "code" and example["concept"] == "Loops"
    assert example["lines"][0] == "for i in range(3):"
    assert "loop" in example["text"].lower()
    lecture._concept_code("Loops", "Python")  # cached -> no second fetch
    assert calls["n"] == 1


def test_attach_concept_examples_one_unit_per_concept():
    import knowledge.lecture as lecture

    # Conditionals/Loops/Functions are all curated, so units come from the
    # library (no network), de-duplicated and attached to the first owner.
    results = [
        lecture.ObjectiveResult(objective="Explain conditionals and loops"),
        lecture.ObjectiveResult(objective="Write functions and more loops"),  # 'loops' dup
    ]
    lecture._attach_concept_examples(results, [r.objective for r in results], "Intro to Python")
    concepts = [unit["concept"] for r in results for unit in r.concept_examples]
    assert concepts == ["Conditionals", "Loops", "Functions"]  # deduped, ordered
    # Each unit is the full 4-part shape.
    for r in results:
        for unit in r.concept_examples:
            assert set(unit) >= {"example", "walkthrough", "practice", "answer"}
    # 'Loops' attached to the FIRST objective that named it
    assert any(unit["concept"] == "Loops" for unit in results[0].concept_examples)
    assert all(unit["concept"] != "Loops" for unit in results[1].concept_examples)


def test_concept_unit_fallback_builds_deterministic_unit(monkeypatch):
    import knowledge.lecture as lecture

    # An uncurated concept: retrieval supplies the example; the rest is derived.
    monkeypatch.setattr(
        lecture,
        "_concept_code",
        lambda concept, language: {
            "kind": "code",
            "concept": concept,
            "text": "A worked example.",
            "lines": ["x = 1", "for i in range(x):", "    print(i)"],
            "title": concept,
            "url": "u",
            "source": "Stack Overflow",
        },
    )
    unit = lecture._concept_unit("Generators", "Python")
    assert unit["example"]["lines"][0] == "x = 1"
    assert len(unit["walkthrough"]) == 3                 # one line described per code line
    assert any("Assigns a value to x" in b for b in unit["walkthrough"])
    assert any("loop" in b.lower() for b in unit["walkthrough"])
    assert unit["practice"] and "Recreate" in unit["practice"][0]


def test_describe_line_recognizes_common_shapes():
    from knowledge.lecture import _describe_line

    assert _describe_line("def greet(name):").startswith("Defines the function greet")
    assert _describe_line("class Dog:").startswith("Defines the class Dog")
    assert "loop" in _describe_line("for i in range(3):").lower()
    assert _describe_line("total = 5 + 2").startswith("Assigns a value to total")
    assert _describe_line("return total").startswith("Returns")
    assert _describe_line("") is None


def test_is_programming_lecture_detection():
    from knowledge.lecture import is_programming_lecture

    assert is_programming_lecture(["write a Python for loop", "define a function"])
    assert is_programming_lecture(["explain the history of it"], title="Introduction to Python")
    assert not is_programming_lecture(
        ["explain operant conditioning", "describe classical conditioning"]
    )


def _slide_code(slide):
    """The monospace code lines on a slide built by these helpers."""
    for shape in slide.shapes:
        if shape.has_text_frame and any(
            p.font.name == "Courier New" for p in shape.text_frame.paragraphs
        ):
            return [p.text for p in shape.text_frame.paragraphs]
    return []


def _unit_example():
    return {
        "concept": "Loops",
        "language": "Python",
        "example": {
            "caption": "For example, this loop prints 0, 1, 2.",
            "lines": ["for i in range(3):", "    print(i)"],
        },
        "walkthrough": ["range(3) yields 0, 1, 2.", "print runs once per value."],
        "practice": ["Print the numbers 0 to 4 with a loop."],
        "answer": {"caption": "Looping to five.", "lines": ["for i in range(5):", "    print(i)"]},
        "title": "",
        "url": "",
        "source": "Curated",
    }


def test_concept_unit_renders_four_slides_in_order():
    from pptx import Presentation

    results = [
        ObjectiveResult(
            objective="Explain loops",
            points=["A loop repeats code."],
            concept_examples=[_unit_example()],
            citations=[],
            confidence="high",
        )
    ]
    deck = Presentation(io.BytesIO(build_module_deck("Intro to Python", results)))
    titles = [slide_title(s) for s in deck.slides]
    unit_titles = [t for t in titles if t.split(":")[0] in ("Example", "Walkthrough", "Practice", "Answer")]
    # The fixed unit, immediately consecutive and in order.
    assert unit_titles == ["Example: Loops", "Walkthrough: Loops", "Practice: Loops", "Answer: Loops"]


def test_concept_example_slide_renders_words_and_code():
    from pptx import Presentation

    deck = Presentation(io.BytesIO(build_module_deck("Intro to Python", [
        ObjectiveResult(objective="Explain loops", points=["A loop repeats code."],
                        concept_examples=[_unit_example()], citations=[], confidence="high")
    ])))
    slide = next(s for s in deck.slides if slide_title(s) == "Example: Loops")
    text = "\n".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
    assert "this loop prints" in text                       # words
    assert "for i in range(3):" in text                     # code
    assert _slide_code(slide)                                # rendered in monospace


def test_practice_code_is_example_reference_not_the_answer():
    # Spec §4: walkthrough and practice both show the EXAMPLE snippet verbatim;
    # only the answer carries its own distinct solution.
    from pptx import Presentation

    deck = Presentation(io.BytesIO(build_module_deck("Intro to Python", [
        ObjectiveResult(objective="Explain loops", points=["A loop repeats code."],
                        concept_examples=[_unit_example()], citations=[], confidence="high")
    ])))
    by_title = {slide_title(s): s for s in deck.slides}
    example_code = _slide_code(by_title["Example: Loops"])
    assert _slide_code(by_title["Walkthrough: Loops"]) == example_code
    assert _slide_code(by_title["Practice: Loops"]) == example_code
    # The answer must NOT be the reference snippet.
    assert _slide_code(by_title["Answer: Loops"]) != example_code


def test_extract_prose_example_by_marker():
    query = analyze("What is cognitive dissonance?")
    passage = make_passage()
    marked = ScoredSentence(
        text="For example, a smoker who knows smoking is harmful feels discomfort.",
        tokens=[],
        passage=passage,
        score=5.0,
    )
    plain = ScoredSentence(text="It is mental discomfort.", tokens=[], passage=passage, score=4.0)
    examples = extract_examples(query, [passage], [marked, plain])
    assert examples[0]["kind"] == "prose"
    assert "For example" in examples[0]["text"]


def test_extract_falls_back_to_top_sentence():
    query = analyze("What is entropy?")
    passage = make_passage()
    plain = ScoredSentence(
        text="Entropy measures disorder in a system over time.",
        tokens=[],
        passage=passage,
        score=3.0,
    )
    examples = extract_examples(query, [passage], [plain])
    assert examples[0]["fallback"] is True


# --- deck assembly -----------------------------------------------------------


def build_two_objective_deck():
    results = [
        ObjectiveResult(
            objective="Define variables",
            points=["A variable stores data.", "Variables have names."],
            examples=[
                {
                    "kind": "prose",
                    "text": "For example, x = 5 binds x to 5.",
                    "title": "Variable",
                    "url": "http://v",
                    "source": "Wikipedia",
                }
            ],
            citations=[{"title": "Variable", "url": "http://v", "source": "Wikipedia"}],
            confidence="high",
        ),
        ObjectiveResult(
            objective="Write loops",
            points=["A loop repeats code."],
            examples=[
                {
                    "kind": "code",
                    "text": "For example, this loop prints 0, 1, 2.",
                    "lines": ["for i in range(3):", "    print(i)"],
                    "title": "Loops",
                    "url": "http://l",
                    "source": "Stack Overflow",
                }
            ],
            citations=[],
            confidence="medium",
        ),
    ]
    return Presentation(io.BytesIO(build_module_deck("Intro to Python", results)))


def test_deck_structure_and_notes():
    deck = build_two_objective_deck()
    titles = [slide_title(s) for s in deck.slides]
    # title, overview, case study, obj1 explanation, obj1 example, obj2 explanation,
    # obj2 example, references
    assert len(deck.slides) == 8
    assert any("Intro to Python" in t for t in titles)
    # Case study is slide 3, right after title + overview, before any concept slides.
    assert titles[2].startswith("Case Study:")
    assert sum(t.startswith("Case Study:") for t in titles) == 1
    assert "Variables" in titles and "Loops" in titles  # clean topic titles, no "Objective N:"
    assert sum(t.startswith("Example") for t in titles) == 2
    assert "References" in titles


def test_deck_examples_have_speaker_notes():
    deck = build_two_objective_deck()
    notes = [
        s.notes_slide.notes_text_frame.text for s in deck.slides if s.has_notes_slide
    ]
    assert any("Talking points" in n for n in notes)
    assert any("Source:" in n for n in notes)


def test_code_example_slide_has_words_and_code():
    deck = build_two_objective_deck()
    example_slide = next(
        s for s in deck.slides
        if any(sh.has_text_frame and "loop prints" in sh.text_frame.text for sh in s.shapes)
    )
    text = "\n".join(sh.text_frame.text for sh in example_slide.shapes if sh.has_text_frame)
    assert "For example, this loop prints" in text          # words
    assert "for i in range(3):" in text                     # code
    # the code lives in a monospace box
    monospace = any(
        sh.has_text_frame
        and any(p.font.name == "Courier New" for p in sh.text_frame.paragraphs)
        for sh in example_slide.shapes
    )
    assert monospace


def test_deck_references_number_unique_sources():
    deck = build_two_objective_deck()
    references = next(s for s in deck.slides if slide_title(s) == "References")
    body = "\n".join(
        shape.text_frame.text for shape in references.shapes if shape.has_text_frame
    )
    assert "Variable" in body and "Loops" in body


# --- API route ---------------------------------------------------------------


@pytest.fixture
def client():
    app_module.app.config["TESTING"] = True
    return app_module.app.test_client()


def test_lecture_route_serves_pptx(client, monkeypatch):
    monkeypatch.delenv("API_KEY", raising=False)
    monkeypatch.setattr(
        service_module,
        "build_lecture_deck",
        lambda objectives, title, source_label=None: (
            b"PK\x03\x04fakepptx", {"objectives": 2, "title": title, "items": []}
        ),
    )
    res = client.post(
        "/api/v1/lecture",
        json={"objectives": "define variables. explain loops.", "title": "T"},
    )
    assert res.status_code == 200
    assert res.mimetype == PPTX_MIMETYPE
    assert res.data == b"PK\x03\x04fakepptx"


def test_lecture_route_accepts_list(client, monkeypatch):
    captured = {}
    monkeypatch.delenv("API_KEY", raising=False)

    def fake(objectives, title, source_label=None):
        captured["objectives"] = objectives
        return (b"x", {"objectives": 2, "title": title, "items": []})

    monkeypatch.setattr(service_module, "build_lecture_deck", fake)
    client.post("/api/v1/lecture", json={"objectives": ["define x", "explain y"]})
    assert "define x" in captured["objectives"] and "explain y" in captured["objectives"]


def test_lecture_route_validates_length(client, monkeypatch):
    monkeypatch.delenv("API_KEY", raising=False)
    res = client.post("/api/v1/lecture", json={"objectives": "short"})
    assert res.status_code == 400
    assert res.get_json()["error"]["code"] == "invalid_request"


def test_lecture_in_openapi(client):
    spec = client.get("/api/v1/openapi.json").get_json()
    assert "/api/v1/lecture" in spec["paths"]
    assert spec["paths"]["/api/v1/lecture"]["post"]["security"] == [{"ApiKeyAuth": []}]


# --- conceptual (non-programming) profile -----------------------------------


def test_classify_subject_programming_vs_conceptual():
    from knowledge.lecture import classify_subject

    assert classify_subject(["write a for loop", "define a function"]) == "programming"
    assert classify_subject(["explain operant conditioning", "describe attachment theory"]) == "conceptual"
    assert classify_subject(["explain photosynthesis"], title="Introduction to Biology") == "conceptual"


def test_review_questions_generation():
    from knowledge.lecture import _review_questions

    qs = _review_questions("cognitive dissonance", "operant conditioning")
    assert qs[0] == "Define cognitive dissonance in your own words."
    assert any("real-world example" in q for q in qs)
    assert qs[-1] == "How does cognitive dissonance relate to operant conditioning?"
    # No next topic -> no relate question.
    assert all("relate to" not in q for q in _review_questions("entropy"))


def test_conceptual_deck_structure():
    results = [
        ObjectiveResult(
            objective="Explain cognitive dissonance",
            points=["It is the mental discomfort of holding conflicting beliefs."],
            examples=[{
                "kind": "prose",
                "text": "For example, a smoker who knows smoking is harmful feels discomfort.",
                "title": "Cognitive dissonance", "url": "u", "source": "Wikipedia",
            }],
            questions=[
                "Define cognitive dissonance in your own words.",
                "Explain why cognitive dissonance is important.",
            ],
            citations=[{"title": "Cognitive dissonance", "url": "u", "source": "Wikipedia"}],
            confidence="high",
        )
    ]
    deck = Presentation(io.BytesIO(build_module_deck("Introduction to Psychology", results)))
    titles = [slide_title(s) for s in deck.slides]
    assert "Cognitive Dissonance" in titles                       # concept slide
    assert any(t.startswith("Illustration:") for t in titles)
    assert any(t.startswith("Check Your Understanding:") for t in titles)
    # No programming code units anywhere in a conceptual deck.
    assert not any(
        t.split(":")[0] in ("Example", "Walkthrough", "Practice", "Answer") for t in titles
    )
    # Case study is the matched field, not the computing default.
    assert any("Stanford Prison Experiment" in t for t in titles)


def test_conceptual_questions_carry_model_answers_in_notes():
    results = [ObjectiveResult(
        objective="Explain entropy",
        points=["Entropy measures disorder in a system."],
        questions=["Define entropy in your own words."],
        citations=[], confidence="medium",
    )]
    deck = Presentation(io.BytesIO(build_module_deck("Thermodynamics", results)))
    cyu = next(s for s in deck.slides if slide_title(s).startswith("Check Your Understanding"))
    notes = cyu.notes_slide.notes_text_frame.text
    assert "Entropy measures disorder in a system." in notes


def test_build_lecture_deck_conceptual_attaches_questions(monkeypatch):
    import knowledge.lecture as lecture

    monkeypatch.setattr(
        lecture,
        "_build_objective",
        lambda objective, context="", programming_lecture=False: ObjectiveResult(
            objective=objective, points=["A relevant point."], confidence="medium"
        ),
    )
    pptx_bytes, summary = lecture.build_lecture_deck(
        "Explain cognitive dissonance. Describe attachment theory.",
        title="Introduction to Psychology",
    )
    titles = [slide_title(s) for s in Presentation(io.BytesIO(pptx_bytes)).slides]
    assert any(t.startswith("Check Your Understanding:") for t in titles)
    assert not any(t.startswith("Example:") for t in titles)


# --- quantitative profile ----------------------------------------------------


def test_classify_subject_quantitative():
    from knowledge.lecture import classify_subject

    assert classify_subject(["solve quadratic equations"], title="Algebra I") == "quantitative"
    assert classify_subject(["apply Newton's second law"], title="Physics") == "quantitative"
    # Programming detection still takes precedence.
    assert classify_subject(["write a function to compute the mean"], title="Python") == "programming"
    # A humanities module stays conceptual.
    assert classify_subject(["discuss the causes of the French Revolution"], title="World History") == "conceptual"


def test_quant_deck_structure_and_no_answer_leak():
    from knowledge.lecture import _attach_quant_units

    results = [ObjectiveResult(objective="Solve quadratic equations",
                               points=["A quadratic equation has the form ax^2 + bx + c = 0."],
                               confidence="high")]
    _attach_quant_units(results)
    assert results[0].quant_units  # matched a curated concept

    deck = Presentation(io.BytesIO(build_module_deck("Algebra I", results)))
    titles = [slide_title(s) for s in deck.slides]
    assert any(t.startswith("Worked Example:") for t in titles)
    assert any(t.startswith("Practice:") for t in titles)
    assert any(t.startswith("Answer:") for t in titles)
    # No code anywhere in a quantitative deck.
    assert not any(
        sh.has_text_frame and any(p.font.name == "Courier New" for p in sh.text_frame.paragraphs)
        for s in deck.slides for sh in s.shapes
    )
    # The Practice slide shows the problem but must NOT reveal the solution.
    practice = next(s for s in deck.slides if slide_title(s).startswith("Practice:"))
    ptext = "\n".join(sh.text_frame.text for sh in practice.shapes if sh.has_text_frame)
    assert "x = 3" not in ptext and "x = 4" not in ptext


def test_quant_unmatched_objective_falls_back_to_conceptual(monkeypatch):
    import knowledge.lecture as lecture

    monkeypatch.setattr(
        lecture,
        "_build_objective",
        lambda objective, context="", programming_lecture=False: ObjectiveResult(
            objective=objective, points=["A point."], confidence="medium"
        ),
    )
    pptx_bytes, _ = lecture.build_lecture_deck(
        "Solve quadratic equations. Discuss the history of mathematics.", title="Algebra"
    )
    titles = [slide_title(s) for s in Presentation(io.BytesIO(pptx_bytes)).slides]
    assert any(t.startswith("Worked Example: Quadratic") for t in titles)   # matched -> quant unit
    assert any(t.startswith("Check Your Understanding:") for t in titles)   # unmatched -> conceptual


# --- file upload on POST /api/v1/lecture ------------------------------------


def _capture_lecture(monkeypatch):
    captured = {}

    def fake(objectives, title, source_label=None):
        captured.update(objectives=objectives, title=title, source_label=source_label)
        return (b"PK\x03\x04fake", {"objectives": 1, "title": title, "items": []})

    monkeypatch.setattr(service_module, "build_lecture_deck", fake)
    return captured


def test_lecture_accepts_text_file_upload(client, monkeypatch):
    monkeypatch.delenv("API_KEY", raising=False)
    captured = _capture_lecture(monkeypatch)
    res = client.post(
        "/api/v1/lecture",
        data={"file": (io.BytesIO(b"Define variables\nExplain loops\nWrite functions"), "notes.txt")},
        content_type="multipart/form-data",
    )
    assert res.status_code == 200
    assert res.mimetype == PPTX_MIMETYPE
    assert "Define variables" in captured["objectives"]
    assert captured["source_label"] == "notes.txt"


def test_lecture_upload_merges_file_then_objectives(client, monkeypatch):
    monkeypatch.delenv("API_KEY", raising=False)
    captured = _capture_lecture(monkeypatch)
    client.post(
        "/api/v1/lecture",
        data={
            "file": (io.BytesIO(b"From the file"), "notes.txt"),
            "objectives": "Typed objective about recursion",
            "title": "My Module",
        },
        content_type="multipart/form-data",
    )
    objectives = captured["objectives"]
    assert objectives.index("From the file") < objectives.index("Typed objective")
    assert captured["title"] == "My Module"


def test_lecture_upload_pptx_extracts_slide_text(client, monkeypatch):
    from pptx import Presentation
    from pptx.util import Inches

    monkeypatch.delenv("API_KEY", raising=False)
    captured = _capture_lecture(monkeypatch)
    deck = Presentation()
    for text in ("Introduction to Supply and Demand", "Market equilibrium"):
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1)).text_frame.text = text
    buf = io.BytesIO()
    deck.save(buf)
    res = client.post(
        "/api/v1/lecture",
        data={"file": (io.BytesIO(buf.getvalue()), "econ.pptx")},
        content_type="multipart/form-data",
    )
    assert res.status_code == 200
    assert "Supply and Demand" in captured["objectives"]


def test_lecture_upload_requires_file_or_objectives(client, monkeypatch):
    monkeypatch.delenv("API_KEY", raising=False)
    res = client.post("/api/v1/lecture", data={}, content_type="multipart/form-data")
    assert res.status_code == 400
    assert res.get_json()["error"]["code"] == "invalid_request"


def test_lecture_upload_unsupported_type(client, monkeypatch):
    monkeypatch.delenv("API_KEY", raising=False)
    res = client.post(
        "/api/v1/lecture",
        data={"file": (io.BytesIO(b"MZ\x90\x00"), "malware.exe")},
        content_type="multipart/form-data",
    )
    assert res.status_code == 415
    assert res.get_json()["error"]["code"] == "unsupported_media_type"


def test_lecture_upload_too_large(client, monkeypatch):
    monkeypatch.delenv("API_KEY", raising=False)
    monkeypatch.setattr(service_module, "MAX_UPLOAD_BYTES", 10)
    res = client.post(
        "/api/v1/lecture",
        data={"file": (io.BytesIO(b"this is more than ten bytes"), "notes.txt")},
        content_type="multipart/form-data",
    )
    assert res.status_code == 413
    assert res.get_json()["error"]["code"] == "payload_too_large"


def test_lecture_upload_no_extractable_text(client, monkeypatch):
    monkeypatch.delenv("API_KEY", raising=False)
    res = client.post(
        "/api/v1/lecture",
        data={"file": (io.BytesIO(b"   "), "empty.txt")},
        content_type="multipart/form-data",
    )
    assert res.status_code == 422
    assert res.get_json()["error"]["code"] == "invalid_request"


def test_lecture_json_path_unchanged(client, monkeypatch):
    monkeypatch.delenv("API_KEY", raising=False)
    _capture_lecture(monkeypatch)
    res = client.post("/api/v1/lecture", json={"objectives": "define variables and explain loops"})
    assert res.status_code == 200


def test_openapi_lecture_documents_multipart(client):
    spec = client.get("/api/v1/openapi.json").get_json()
    content = spec["paths"]["/api/v1/lecture"]["post"]["requestBody"]["content"]
    assert "multipart/form-data" in content and "application/json" in content
